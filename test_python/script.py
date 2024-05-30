import csv
from pptx import Presentation
from pptx.util import Inches

def load_csv(csv_file):
    data = {}
    with open(csv_file, mode='r', encoding='utf-8') as file:
        reader = csv.DictReader(file)
        for row in reader:
            data[row['Key']] = row['Data']
    return data

def verify_pptx(pptx_file, keys):
    prs = Presentation(pptx_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'alt_text') and shape.alt_text:
                if shape.alt_text not in keys:
                    print(f"Key {shape.alt_text} not found in CSV.")
                    return False
    return True

def update_pptx(pptx_file, data):
    prs = Presentation(pptx_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'alt_text') and shape.alt_text:
                key = shape.alt_text
                if key in data:
                    if shape.has_text_frame:
                        shape.text = data[key]
                    elif shape.shape_type == 13:  # Picture shape type
                        image_path = data[key]
                        shape._element.clear_content()
                        shape._element.getparent().remove(shape._element)
                        slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, shape.height)
    prs.save('updated_presentation.pptx')

# Load CSV data
csv_file = 'data.csv'
data = load_csv(csv_file)

# Verify and update PPTX
pptx_file = 'presentation.pptx'
if verify_pptx(pptx_file, data.keys()):
    update_pptx(pptx_file, data)
else:
    print("Verification failed.")
